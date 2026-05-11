"""
DocBook - Doctor Booking App
Generates PowerPoint presentations for each feature/functionality.
Run: python generate_presentations.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.util as util

# ── Color palette ──────────────────────────────────────────────────────────────
PRIMARY       = RGBColor(0x1A, 0x73, 0xE8)   # #1A73E8
PRIMARY_DARK  = RGBColor(0x15, 0x57, 0xB0)   # #1557B0
SECONDARY     = RGBColor(0x34, 0xA8, 0x53)   # #34A853
ACCENT        = RGBColor(0xFF, 0x6B, 0x6B)   # #FF6B6B
BACKGROUND    = RGBColor(0xF8, 0xF9, 0xFA)   # #F8F9FA
SURFACE       = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
TEXT_PRIMARY  = RGBColor(0x1A, 0x1A, 0x2E)   # #1A1A2E
TEXT_SEC      = RGBColor(0x6B, 0x72, 0x80)   # #6B7280
STAR_GOLD     = RGBColor(0xFB, 0xBC, 0x04)   # #FBBC04
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # dark slide backgrounds

OUT_DIR = "presentations"
os.makedirs(OUT_DIR, exist_ok=True)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO if False else 1,  # MSO_SHAPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=TEXT_PRIMARY,
                align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_title_slide(prs, title, subtitle, bg_color=PRIMARY, accent_color=None):
    slide = blank_slide(prs)
    fill_bg(slide, bg_color)

    # Left accent bar
    bar = add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H,
                   accent_color or SECONDARY)

    # Logo / app name top-right
    add_textbox(slide, "DocBook", Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.6),
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # Main title
    add_textbox(slide, title, Inches(0.8), Inches(2.2), Inches(9), Inches(1.4),
                font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, subtitle, Inches(0.8), Inches(3.8), Inches(9), Inches(0.8),
                font_size=24, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF),
                align=PP_ALIGN.LEFT)

    # Bottom bar
    add_rect(slide, Inches(0), Inches(6.9), SLIDE_W, Inches(0.6),
             accent_color or SECONDARY)
    add_textbox(slide, "DocBook — Doctor Booking & Appointment App",
                Inches(0.3), Inches(6.85), Inches(9), Inches(0.5),
                font_size=13, color=WHITE)
    return slide


def add_section_slide(prs, heading, items, icon_char="▶",
                       left_color=PRIMARY, bg=BACKGROUND):
    slide = blank_slide(prs)
    fill_bg(slide, bg)

    # Header strip
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.4), left_color)
    add_textbox(slide, f"{icon_char}  {heading}",
                Inches(0.4), Inches(0.3), Inches(12), Inches(0.85),
                font_size=34, bold=True, color=WHITE)

    # Feature bullets
    y = Inches(1.7)
    for item in items:
        bullet_dot = add_rect(slide, Inches(0.4), y + Inches(0.08),
                               Inches(0.15), Inches(0.15), left_color)
        add_textbox(slide, item, Inches(0.75), y, Inches(11.8), Inches(0.45),
                    font_size=17, color=TEXT_PRIMARY)
        y += Inches(0.52)
        if y > Inches(6.8):
            break
    return slide


def add_two_col_slide(prs, heading, left_items, right_items,
                      left_header="Features", right_header="Details",
                      bg=BACKGROUND, hdr_color=PRIMARY):
    slide = blank_slide(prs)
    fill_bg(slide, bg)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), hdr_color)
    add_textbox(slide, heading, Inches(0.4), Inches(0.2), Inches(12), Inches(0.8),
                font_size=32, bold=True, color=WHITE)

    # Left column header
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(5.8), Inches(0.5), hdr_color)
    add_textbox(slide, left_header, Inches(0.5), Inches(1.38), Inches(5.6), Inches(0.45),
                font_size=15, bold=True, color=WHITE)

    # Right column header
    add_rect(slide, Inches(6.8), Inches(1.35), Inches(6.1), Inches(0.5), hdr_color)
    add_textbox(slide, right_header, Inches(6.9), Inches(1.38), Inches(5.9), Inches(0.45),
                font_size=15, bold=True, color=WHITE)

    # Left items
    y = Inches(2.05)
    for item in left_items:
        add_textbox(slide, f"• {item}", Inches(0.5), y, Inches(5.6), Inches(0.5),
                    font_size=15, color=TEXT_PRIMARY)
        y += Inches(0.48)
        if y > Inches(6.8): break

    # Right items
    y = Inches(2.05)
    for item in right_items:
        add_textbox(slide, f"→ {item}", Inches(6.9), y, Inches(5.8), Inches(0.5),
                    font_size=15, color=TEXT_PRIMARY)
        y += Inches(0.48)
        if y > Inches(6.8): break

    return slide


def add_screenshot_placeholder(prs, label, caption, phone_color=PRIMARY):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)

    # Phone frame
    phone_l, phone_t = Inches(4.4), Inches(0.4)
    phone_w, phone_h = Inches(4.5), Inches(6.7)

    phone_shape = add_rect(slide, phone_l, phone_t, phone_w, phone_h,
                           RGBColor(0x22, 0x22, 0x33),
                           RGBColor(0x44, 0x55, 0x77))

    # Screen area inside phone
    scr_l = phone_l + Inches(0.25)
    scr_t = phone_t + Inches(0.55)
    scr_w = phone_w - Inches(0.5)
    scr_h = phone_h - Inches(1.0)
    add_rect(slide, scr_l, scr_t, scr_w, scr_h, phone_color)

    # Screen label
    add_textbox(slide, label, scr_l + Inches(0.2), scr_t + Inches(0.5),
                scr_w - Inches(0.4), Inches(1.2),
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Caption below phone
    add_textbox(slide, caption,
                Inches(1), Inches(6.85), Inches(11.3), Inches(0.55),
                font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                align=PP_ALIGN.CENTER)

    # Screenshot note
    add_textbox(slide,
                "[ Screenshot: Place actual screenshot here — assets/screenshots/ ]",
                scr_l + Inches(0.1), scr_t + scr_h - Inches(0.9),
                scr_w - Inches(0.2), Inches(0.7),
                font_size=9, color=RGBColor(0xDD, 0xDD, 0xFF),
                align=PP_ALIGN.CENTER)

    return slide


def add_flow_slide(prs, heading, steps, bg=BACKGROUND, hdr_color=PRIMARY):
    slide = blank_slide(prs)
    fill_bg(slide, bg)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), hdr_color)
    add_textbox(slide, heading, Inches(0.4), Inches(0.2), Inches(12), Inches(0.8),
                font_size=32, bold=True, color=WHITE)

    box_w = Inches(1.9)
    box_h = Inches(1.4)
    gap   = Inches(0.2)
    start_x = Inches(0.5)
    y_top = Inches(2.0)

    for i, (step_num, step_title, step_desc) in enumerate(steps):
        x = start_x + i * (box_w + gap + Inches(0.5))

        # Step box
        add_rect(slide, x, y_top, box_w, box_h, hdr_color)
        add_textbox(slide, step_num, x, y_top + Inches(0.05), box_w, Inches(0.5),
                    font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, step_title, x, y_top + Inches(0.55), box_w, Inches(0.7),
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow
        if i < len(steps) - 1:
            arr_x = x + box_w + Inches(0.1)
            add_textbox(slide, "→", arr_x, y_top + Inches(0.5), Inches(0.4), Inches(0.5),
                        font_size=22, bold=True, color=hdr_color, align=PP_ALIGN.CENTER)

        # Description below box
        add_textbox(slide, step_desc, x - Inches(0.1), y_top + box_h + Inches(0.2),
                    box_w + Inches(0.2), Inches(1.5),
                    font_size=12, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 1 — Home Screen / Doctor Discovery
# ══════════════════════════════════════════════════════════════════════════════

def make_p01():
    prs = new_prs()

    add_title_slide(prs,
        "Doctor Discovery\n& Search",
        "Home Screen — Browse, Filter & Find Specialists",
        bg_color=PRIMARY, accent_color=SECONDARY)

    add_section_slide(prs, "Overview", [
        "Entry point of the app — displays all 37 available doctors",
        "Real-time search by doctor name, specialty, or hospital",
        "Multi-filter system: City × Specialty × Rating",
        "Live result count updates as filters change",
        "Empty state with 'Clear Filters' recovery shortcut",
        "Tap any card to open the full Doctor Detail profile",
        "Quick access to My Bookings via top-right calendar icon",
    ], icon_char="🏠", left_color=PRIMARY)

    add_two_col_slide(prs,
        "Search & Filter Features",
        ["Search Bar", "City Filter", "Specialty Filter",
         "Rating Filter", "Clear All Button", "Results Count"],
        ["Searches name, specialty & hospital in real time",
         "8 cities: Mumbai, Delhi, Bangalore, Chennai, Hyderabad, Pune, Kolkata, Ahmedabad",
         "8 specialties: Cardiology, Dermatology, Pediatrics, Orthopedic, Gynecology, Neurology, Psychiatry, Ophthalmology",
         "Options: All / 3★+ / 3.5★+ / 4★+ / 4.5★+",
         "Resets all filters and shows full list",
         "Live count of matching doctors"],
        left_header="Filter Type", right_header="Behaviour",
        hdr_color=PRIMARY)

    add_section_slide(prs, "Doctor Card — Information Shown", [
        "Color-coded avatar with doctor initials",
        "Full name and medical specialty",
        "Years of clinical experience",
        "Star rating (4.4 – 4.9) with patient review count",
        "Hospital name and city",
        "Consultation fee in ₹ (₹500 – ₹1,500)",
        "Tap card → opens Doctor Detail Screen",
    ], icon_char="🩺", left_color=PRIMARY)

    add_screenshot_placeholder(prs,
        "Home Screen\nDoctor List",
        "assets/screenshots/home_screen.png  |  assets/screenshots/home_filtered.png",
        phone_color=PRIMARY)

    add_flow_slide(prs, "User Journey from Home Screen",
        [("1", "Open App", "Home screen\nloads 37 doctors"),
         ("2", "Search / Filter", "Apply city,\nspecialty, rating"),
         ("3", "Browse Results", "Live-updated\ndoctor cards"),
         ("4", "Select Doctor", "Tap card to\nopen profile"),
         ("5", "Book", "Proceed to\nappointment")],
        hdr_color=PRIMARY)

    prs.save(os.path.join(OUT_DIR, "01_home_screen.pptx"))
    print("[OK] 01_home_screen.pptx")


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 2 — Doctor Detail Screen
# ══════════════════════════════════════════════════════════════════════════════

def make_p02():
    prs = new_prs()

    add_title_slide(prs,
        "Doctor Profile\nDetail Screen",
        "Full Doctor Information — Qualifications, Slots & Booking CTA",
        bg_color=RGBColor(0x10, 0x4E, 0xA8), accent_color=STAR_GOLD)

    add_section_slide(prs, "What's Displayed on the Profile", [
        "Gradient header with large color-coded initials avatar",
        "Full name, medical specialty, and qualifications (e.g., MBBS, MD, DM)",
        "Star rating (out of 5) with total patient review count",
        "Hospital name with map pin icon",
        "Three-stat row: Years of experience · Patients served · Consultation fee",
        "Detailed 'About' bio describing expertise and approach",
        "Available working days shown as interactive chips",
        "Time slot grid for the selected day",
        "Consultation type options: In-Person / Video / Phone",
        "Book Appointment CTA button",
    ], icon_char="👨‍⚕️", left_color=RGBColor(0x10, 0x4E, 0xA8))

    add_two_col_slide(prs,
        "Doctor Data Points",
        ["Qualifications", "Rating", "Experience", "Hospital",
         "Consultation Fee", "Available Days", "Time Slots", "Consultation Types"],
        ["e.g., MBBS, MD (Cardiology), DM",
         "4.4 – 4.9 ⭐ with 98–567 patient reviews",
         "8 – 25 years of clinical experience",
         "Named real hospital in chosen city",
         "₹500 – ₹1,500 base fee",
         "Mon – Sat patterns (varies per doctor)",
         "4–6 slots per day (9:00 AM – 6:00 PM range)",
         "In-Person · Video Call · Phone Call"],
        left_header="Data Field", right_header="Values / Range",
        hdr_color=RGBColor(0x10, 0x4E, 0xA8))

    add_section_slide(prs, "Consultation Types & Pricing Logic", [
        "In-Person Consultation → 100% of base fee  (e.g., ₹900)",
        "Video Call Consultation → 80% of base fee  (e.g., ₹720)",
        "Phone Call Consultation → 60% of base fee  (e.g., ₹540)",
        "Pricing is dynamically calculated on the Booking Screen",
        "All types include 18% GST added at checkout",
        "Doctor's available days restrict the date picker in Booking",
        "Patient can review all info before committing to Book",
    ], icon_char="💼", left_color=RGBColor(0x10, 0x4E, 0xA8))

    add_screenshot_placeholder(prs,
        "Doctor\nDetail Screen",
        "assets/screenshots/doctor_detail.png",
        phone_color=RGBColor(0x10, 0x4E, 0xA8))

    prs.save(os.path.join(OUT_DIR, "02_doctor_detail.pptx"))
    print("[OK] 02_doctor_detail.pptx")


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 3 — Booking Screen
# ══════════════════════════════════════════════════════════════════════════════

def make_p03():
    prs = new_prs()

    add_title_slide(prs,
        "Appointment\nBooking Screen",
        "Choose Consultation Type · Date · Time Slot · Proceed to Payment",
        bg_color=RGBColor(0x00, 0x7A, 0x5E), accent_color=STAR_GOLD)

    add_section_slide(prs, "Booking Screen Sections", [
        "Doctor summary bar: Name, specialty & hospital for quick reference",
        "Consultation Type selector (In-Person / Video / Phone) with dynamic pricing",
        "Horizontal date picker showing next 14 available dates",
        "Time slot grid — only shown once a date is selected",
        "Fee Breakdown panel: Base fee + 18% GST = Total",
        "'Proceed to Payment' button — enabled only when date & time are chosen",
    ], icon_char="📅", left_color=RGBColor(0x00, 0x7A, 0x5E))

    add_two_col_slide(prs,
        "Consultation Types & Dynamic Pricing",
        ["In-Person", "Video Call", "Phone Call",
         "GST", "Total", "Button State"],
        ["Full base fee (100%) — e.g., ₹900",
         "80% of base fee — e.g., ₹720",
         "60% of base fee — e.g., ₹540",
         "18% applied on selected fee",
         "Fee + GST shown live in breakdown panel",
         "Disabled until date AND time are both selected"],
        left_header="Option", right_header="Pricing Logic",
        hdr_color=RGBColor(0x00, 0x7A, 0x5E))

    add_section_slide(prs, "Date & Time Slot Selection", [
        "Date picker is a horizontally scrollable row of chips",
        "Shows next 14 dates starting from today",
        "Only dates matching doctor's available days are shown",
        "Each chip displays: Day abbreviation (Mon), Date number, Month",
        "Selected date is highlighted in primary blue",
        "Time slot grid appears below after a date is tapped",
        "Slots show times like 09:00 AM, 10:00 AM … 05:00 PM",
        "Selected slot is highlighted; only one slot can be chosen at a time",
    ], icon_char="🕐", left_color=RGBColor(0x00, 0x7A, 0x5E))

    add_screenshot_placeholder(prs,
        "Booking Screen",
        "assets/screenshots/booking_screen.png",
        phone_color=RGBColor(0x00, 0x7A, 0x5E))

    add_flow_slide(prs, "Booking Flow Steps",
        [("1", "Choose Type", "In-Person /\nVideo / Phone"),
         ("2", "Pick Date", "Next 14 dates\nscroll left-right"),
         ("3", "Pick Time", "Grid of\ntime slots"),
         ("4", "Review Fee", "Base + GST\nbreakdown"),
         ("5", "Proceed", "Go to\nPayment Screen")],
        hdr_color=RGBColor(0x00, 0x7A, 0x5E))

    prs.save(os.path.join(OUT_DIR, "03_booking_screen.pptx"))
    print("[OK] 03_booking_screen.pptx")


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 4 — Payment Screen
# ══════════════════════════════════════════════════════════════════════════════

def make_p04():
    prs = new_prs()

    add_title_slide(prs,
        "Payment Screen\n4 Payment Methods",
        "UPI · Card · Wallet · Net Banking — Secure Checkout",
        bg_color=RGBColor(0x5F, 0x25, 0x9F), accent_color=STAR_GOLD)

    add_section_slide(prs, "Payment Screen Overview", [
        "Order summary card (gradient): Doctor, type, date/time, total with GST",
        "'Secure' badge with lock icon in app bar",
        "4 payment method tabs: UPI | Card | Wallet | Net Banking",
        "Dynamic form area changes based on selected payment method",
        "Pay button shows total amount: 'Pay ₹X Securely'",
        "2-second processing animation with spinner on pay tap",
        "On success → navigates to Confirmation Screen",
    ], icon_char="💳", left_color=RGBColor(0x5F, 0x25, 0x9F))

    add_two_col_slide(prs,
        "Payment Method 1 — UPI",
        ["QR Code Display", "UPI ID Entry", "Supported Apps", "UPI Address"],
        ["Custom-painted QR code in app (no external image)",
         "Text field: yourname@upi format",
         "Google Pay, PhonePe, Paytm, BHIM UPI",
         "docbook@ybl (shown below QR)"],
        left_header="Element", right_header="Detail",
        hdr_color=RGBColor(0x5F, 0x25, 0x9F))

    add_two_col_slide(prs,
        "Payment Method 2 — Card",
        ["Saved Cards", "Card Preview", "Card Number", "Holder Name", "Expiry + CVV"],
        ["2 pre-saved: Visa •••• 4242 · Mastercard •••• 8888",
         "Live gradient card preview updates as user types",
         "16-digit field with placeholder •••• •••• •••• ••••",
         "Name on card (auto-capitalised in preview)",
         "MM/YY + obscured CVV field"],
        left_header="Field", right_header="Behaviour",
        hdr_color=RGBColor(0x5F, 0x25, 0x9F))

    add_two_col_slide(prs,
        "Payment Method 3 & 4 — Wallet & Net Banking",
        ["Paytm Wallet", "PhonePe Wallet", "Google Pay Wallet",
         "Amazon Pay Wallet", "Net Banking — 6 Banks"],
        ["Color: #00BAF2 — wallet icon",
         "Color: #5F259F — phone icon",
         "Color: #4285F4 — G icon",
         "Color: #FF9900 — shopping bag icon",
         "SBI · HDFC · ICICI · Axis · Kotak · Bank of Baroda"],
        left_header="Wallet", right_header="Bank / Detail",
        hdr_color=RGBColor(0x5F, 0x25, 0x9F))

    add_screenshot_placeholder(prs,
        "Payment Screen\n(UPI | Card | Wallet | Net Banking)",
        "assets/screenshots/payment_upi.png  |  payment_card.png  |  payment_wallet.png  |  payment_netbanking.png",
        phone_color=RGBColor(0x5F, 0x25, 0x9F))

    prs.save(os.path.join(OUT_DIR, "04_payment_screen.pptx"))
    print("[OK] 04_payment_screen.pptx")


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 5 — Confirmation Screen
# ══════════════════════════════════════════════════════════════════════════════

def make_p05():
    prs = new_prs()

    add_title_slide(prs,
        "Booking\nConfirmation",
        "Success Animation · Booking Summary · Navigation Options",
        bg_color=SECONDARY, accent_color=PRIMARY)

    add_section_slide(prs, "Confirmation Screen Features", [
        "Animated checkmark icon — scale + fade-in animation on load",
        "Large 'Booking Confirmed!' heading in green",
        "Unique 8-character alphanumeric Booking ID",
        "Doctor name, specialty, and hospital",
        "Appointment date (formatted: Mon, Jan 15)",
        "Appointment time slot (e.g., 10:00 AM)",
        "Consultation type (In-Person / Video Call / Phone Call)",
        "Amount paid (including GST breakdown)",
        "Payment method and transaction ID",
        "Two action buttons: Back to Home · View My Bookings",
    ], icon_char="✅", left_color=SECONDARY)

    add_two_col_slide(prs,
        "Booking ID & Transaction Details",
        ["Booking ID", "Transaction ID", "Status", "Date Format",
         "Amount", "Payment Method"],
        ["8-char UUID segment — e.g., A3F7BC12",
         "TXN + Unix timestamp millis — e.g., TXN1715123456789",
         "Confirmed (green badge)",
         "EEE, MMM d — e.g., Mon, Jan 15",
         "Total with 18% GST included",
         "UPI / Card / Wallet / Net Banking label"],
        left_header="Field", right_header="Format / Example",
        hdr_color=SECONDARY)

    add_screenshot_placeholder(prs,
        "Confirmation Screen",
        "assets/screenshots/confirmation_screen.png",
        phone_color=SECONDARY)

    prs.save(os.path.join(OUT_DIR, "05_confirmation_screen.pptx"))
    print("[OK] 05_confirmation_screen.pptx")


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 6 — My Bookings Screen
# ══════════════════════════════════════════════════════════════════════════════

def make_p06():
    prs = new_prs()

    add_title_slide(prs,
        "My Bookings\nDashboard",
        "View & Manage All Appointments — Status Tracking",
        bg_color=RGBColor(0xE6, 0x5C, 0x00), accent_color=PRIMARY)

    add_section_slide(prs, "My Bookings Screen Features", [
        "Lists all booked appointments in reverse-chronological order",
        "Each booking card shows full appointment at a glance",
        "Color-coded status badge: Confirmed (green) · Pending (yellow) · Cancelled (red) · Completed (gray)",
        "Doctor name, specialty, and hospital on each card",
        "Consultation type shown with icon",
        "Formatted appointment date and time",
        "Amount paid on the card",
        "Empty state when no bookings exist with link to Browse Doctors",
    ], icon_char="📋", left_color=RGBColor(0xE6, 0x5C, 0x00))

    add_two_col_slide(prs,
        "Booking Card Fields",
        ["Booking ID", "Status", "Doctor Info", "Consultation Type",
         "Date", "Time", "Amount"],
        ["Unique 8-char ID in card header",
         "Confirmed / Pending / Cancelled / Completed",
         "Name, specialty, hospital",
         "Icon + label: In-Person / Video / Phone",
         "Formatted date: Mon, Jan 15",
         "e.g., 10:00 AM",
         "₹X (total including GST)"],
        left_header="Field", right_header="Values",
        hdr_color=RGBColor(0xE6, 0x5C, 0x00))

    add_two_col_slide(prs,
        "Booking Status Types",
        ["Confirmed", "Pending", "Cancelled", "Completed"],
        ["Green badge — payment received, appointment scheduled",
         "Yellow badge — awaiting confirmation",
         "Red badge — appointment was cancelled",
         "Gray badge — appointment completed"],
        left_header="Status", right_header="Meaning & Badge Color",
        hdr_color=RGBColor(0xE6, 0x5C, 0x00))

    add_screenshot_placeholder(prs,
        "My Bookings\nDashboard",
        "assets/screenshots/my_bookings.png  |  assets/screenshots/my_bookings_empty.png",
        phone_color=RGBColor(0xE6, 0x5C, 0x00))

    prs.save(os.path.join(OUT_DIR, "06_my_bookings.pptx"))
    print("[OK] 06_my_bookings.pptx")


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 7 — Full App Overview
# ══════════════════════════════════════════════════════════════════════════════

def make_p07():
    prs = new_prs()

    add_title_slide(prs,
        "DocBook\nComplete Overview",
        "Doctor Booking & Appointment App — All Features Summary",
        bg_color=DARK_BG, accent_color=PRIMARY)

    add_two_col_slide(prs,
        "App Feature Matrix",
        ["Home / Discovery", "Doctor Detail", "Booking Screen",
         "Payment Screen", "Confirmation", "My Bookings"],
        ["Search + 3-axis filter, 37 doctors, 8 cities",
         "Profile, qualifications, time slots, consultation types",
         "Type select, date picker, time grid, GST breakdown",
         "UPI (QR+ID), Card (saved+new), Wallet, Net Banking",
         "Animated success, booking summary, transaction ID",
         "Status-tracked dashboard with 4 booking states"],
        left_header="Screen", right_header="Key Features",
        hdr_color=PRIMARY)

    add_flow_slide(prs, "Complete User Flow",
        [("1", "Discover", "Browse & filter\n37 doctors"),
         ("2", "Profile", "View doctor\ndetails & slots"),
         ("3", "Book", "Date, time &\nconsult type"),
         ("4", "Pay", "UPI / Card /\nWallet / Bank"),
         ("5", "Confirm", "Booking ID\n& summary"),
         ("6", "Manage", "My Bookings\ndashboard")],
        hdr_color=PRIMARY)

    add_two_col_slide(prs,
        "Technology & Data Highlights",
        ["Flutter 3 + Dart", "Material Design 3", "37 Doctors", "8 Cities",
         "8 Specialties", "4 Payment Methods"],
        ["Cross-platform mobile (iOS + Android + Web)",
         "Google Fonts Poppins, animated containers, custom painters",
         "Ratings 4.4–4.9★, fees ₹500–₹1500, experience 8–25 yrs",
         "Mumbai, Delhi, Bangalore, Chennai, Hyderabad, Pune, Kolkata, Ahmedabad",
         "Cardiology, Dermatology, Pediatrics, Orthopedic, Gynecology, Neurology, Psychiatry, Ophthalmology",
         "UPI + QR Code, Debit/Credit Card, Digital Wallets, Net Banking"],
        left_header="Technology", right_header="Details",
        hdr_color=PRIMARY)

    add_section_slide(prs, "Color Design System", [
        "Primary Blue #1A73E8 — app bars, buttons, selected chips",
        "Primary Dark #1557B0 — gradients, pressed/hover states",
        "Secondary Green #34A853 — success, pay button, confirmation",
        "Accent Red #FF6B6B — alerts, cancel actions",
        "Background #F8F9FA — screen backgrounds",
        "Surface White #FFFFFF — cards, form containers",
        "Text Primary #1A1A2E — headings and body copy",
        "Text Secondary #6B7280 — subtitles, hints, labels",
        "Star Gold #FBBC04 — star rating indicators",
    ], icon_char="🎨", left_color=PRIMARY)

    prs.save(os.path.join(OUT_DIR, "07_app_overview.pptx"))
    print("[OK] 07_app_overview.pptx")


# ── Main ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print(f"\nGenerating DocBook PowerPoint presentations -> {OUT_DIR}/\n")
    make_p01()
    make_p02()
    make_p03()
    make_p04()
    make_p05()
    make_p06()
    make_p07()
    print(f"\nAll 7 presentations saved to: {os.path.abspath(OUT_DIR)}/")
    print("\nFiles generated:")
    for f in sorted(os.listdir(OUT_DIR)):
        if f.endswith(".pptx"):
            size = os.path.getsize(os.path.join(OUT_DIR, f))
            print(f"  {f}  ({size:,} bytes)")
